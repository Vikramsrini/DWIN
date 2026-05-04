from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation('Template for Maintenance Process Model.pptx')

for i, slide in enumerate(prs.slides, 1):
    print(f'\n=== SLIDE {i} ===')
    print(f'Slide width: {prs.slide_width.inches:.2f}in, height: {prs.slide_height.inches:.2f}in')
    print(f'Layout name: {slide.slide_layout.name}')
    
    for j, shape in enumerate(slide.shapes, 1):
        print(f'\n  Shape {j}:')
        print(f'    Type: {shape.shape_type}')
        print(f'    Name: {shape.name}')
        if hasattr(shape, 'left') and hasattr(shape, 'top'):
            print(f'    Position: left={shape.left.inches:.2f}in, top={shape.top.inches:.2f}in')
        if hasattr(shape, 'width') and hasattr(shape, 'height'):
            print(f'    Size: width={shape.width.inches:.2f}in, height={shape.height.inches:.2f}in')
        if hasattr(shape, 'text'):
            text = shape.text.strip()
            if text:
                print(f'    Text: {text[:200]}')
        if shape.has_table:
            table = shape.table
            print(f'    Table: {len(table.rows)} rows x {len(table.columns)} cols')
            for r in range(len(table.rows)):
                row = table.rows[r]
                cells_text = [cell.text[:40] for cell in row.cells]
                print(f'      Row {r}: {cells_text}')
        if shape.has_chart:
            chart = shape.chart
            print(f'    Chart type: {chart.chart_type}')
        if shape.is_placeholder:
            print(f'    Placeholder idx: {shape.placeholder_format.idx}, type: {shape.placeholder_format.type}')
