"""
Generate polished Maintenance Process Model PPTX
Project: Digital Twin Health Activity Tracker
Model:   Iterative Enhancement Model
"""
import os
import math
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import matplotlib.patheffects as pe
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

# ── Constants ────────────────────────────────────────────────────────────────
PROJECT  = "Digital Twin Health Activity Tracker"
DATE     = "29 April 2026"
MODEL    = "Iterative Enhancement Model"
MEMBERS  = "Project Development Team"

W, H = Inches(13.33), Inches(7.5)

# PPTX RGBColor palette
C_DARK   = RGBColor(0x1F, 0x3A, 0x5F)   # dark navy
C_MID    = RGBColor(0x2E, 0x6D, 0xA8)   # medium blue
C_LIGHT  = RGBColor(0x4A, 0x90, 0xD9)   # sky blue
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY   = RGBColor(0xF4, 0xF7, 0xFB)   # slide bg
C_GREEN  = RGBColor(0x1B, 0x7A, 0x3B)
C_ORANGE = RGBColor(0xD9, 0x6B, 0x00)
C_RED    = RGBColor(0xB7, 0x1C, 0x1C)
C_PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
C_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
C_SUB    = RGBColor(0x4A, 0x4A, 0x6A)

# Matplotlib palette (0-1 floats)
M_DARK   = (0x1F/255, 0x3A/255, 0x5F/255)
M_MID    = (0x2E/255, 0x6D/255, 0xA8/255)
M_LIGHT  = (0x4A/255, 0x90/255, 0xD9/255)
M_GREEN  = (0x1B/255, 0x7A/255, 0x3B/255)
M_ORANGE = (0xD9/255, 0x6B/255, 0x00/255)
M_RED    = (0xB7/255, 0x1C/255, 0x1C/255)
M_PURPLE = (0x6A/255, 0x1B/255, 0x9A/255)
M_GREY   = (0.93, 0.95, 0.98)


# ── Helpers ──────────────────────────────────────────────────────────────────

def set_bg(slide, hex_color):
    """Solid background fill for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def add_rect(slide, left, top, w, h, fill_color, line_color=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_tb(slide, left, top, w, h, text, size=14, bold=False, color=C_TEXT,
           align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = 'Calibri'
    return tb

def add_divider(slide, top, color=C_MID, left=Inches(0.8), width=Inches(11.73)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()

BANNER_H = Inches(1.45)

def slide_header(slide, title_text, subtitle_text=None):
    """Standard header: dark banner at top."""
    add_rect(slide, Inches(0), Inches(0), W, BANNER_H, C_DARK)
    if subtitle_text:
        add_tb(slide, Inches(0.5), Inches(0.1), Inches(12.5), Inches(0.62),
               title_text, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
        add_tb(slide, Inches(0.5), Inches(0.82), Inches(12.5), Inches(0.38),
               subtitle_text, size=13, bold=False, color=RGBColor(0xB0, 0xC8, 0xF0),
               align=PP_ALIGN.LEFT)
    else:
        add_tb(slide, Inches(0.5), Inches(0.38), Inches(12.5), Inches(0.62),
               title_text, size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # thin accent line under banner
    add_rect(slide, Inches(0), BANNER_H, W, Pt(4), C_LIGHT)


# ── Diagram generators ───────────────────────────────────────────────────────

def make_iterative_cycle_diagram():
    """Circular arrows showing the 6-step iterative maintenance cycle."""
    fig, ax = plt.subplots(figsize=(10, 7), facecolor='white')
    ax.set_aspect('equal'); ax.axis('off')
    ax.set_xlim(-1.2, 1.2); ax.set_ylim(-1.2, 1.35)

    steps = [
        ("Identify\nChange", M_DARK),
        ("Analyze\nImpact", M_MID),
        ("Design\nSolution", M_LIGHT),
        ("Implement\n& Code", M_GREEN),
        ("Test &\nValidate", M_ORANGE),
        ("Deploy &\nMonitor", M_RED),
    ]
    n = len(steps)
    r_circle = 0.72
    r_box    = 0.18

    for i, (label, color) in enumerate(steps):
        angle = math.pi / 2 - i * 2 * math.pi / n
        cx = r_circle * math.cos(angle)
        cy = r_circle * math.sin(angle)

        # Draw rounded box
        box = FancyBboxPatch((cx - r_box, cy - r_box),
                             2*r_box, 2*r_box,
                             boxstyle="round,pad=0.04",
                             facecolor=color, edgecolor='white', linewidth=2,
                             zorder=3)
        ax.add_patch(box)
        ax.text(cx, cy, label, ha='center', va='center',
                fontsize=9.5, fontweight='bold', color='white', zorder=4,
                linespacing=1.3)

        # Draw arc arrow to next step
        a_next = math.pi / 2 - (i + 1) * 2 * math.pi / n
        start_x = r_circle * math.cos(angle)
        start_y = r_circle * math.sin(angle)
        end_x   = r_circle * math.cos(a_next)
        end_y   = r_circle * math.sin(a_next)

        # offset arrow start/end to edge of boxes
        dx = end_x - start_x; dy = end_y - start_y
        dist = math.sqrt(dx**2 + dy**2)
        sx = start_x + (r_box + 0.04) * dx/dist
        sy = start_y + (r_box + 0.04) * dy/dist
        ex = end_x   - (r_box + 0.05) * dx/dist
        ey = end_y   - (r_box + 0.05) * dy/dist

        ax.annotate('', xy=(ex, ey), xytext=(sx, sy),
                    arrowprops=dict(arrowstyle='->', lw=2.2,
                                   color='#555555',
                                   connectionstyle='arc3,rad=0.22'))

    # Centre label
    centre_circle = plt.Circle((0, 0), 0.26, color=M_DARK, zorder=2)
    ax.add_patch(centre_circle)
    ax.text(0, 0.04, 'Iterative\nEnhancement', ha='center', va='center',
            fontsize=9, fontweight='bold', color='white', zorder=5)

    ax.set_title('Iterative Enhancement Maintenance Cycle',
                 fontsize=14, fontweight='bold', color=tuple(M_DARK), pad=12)
    plt.tight_layout()
    plt.savefig('img_cycle.png', dpi=160, bbox_inches='tight', facecolor='white')
    plt.close()


def make_maintenance_types_diagram():
    """2×2 grid of Boehm maintenance types."""
    fig, axes = plt.subplots(2, 2, figsize=(11, 6.5), facecolor='white')
    fig.subplots_adjust(hspace=0.35, wspace=0.25)

    types = [
        ("CORRECTIVE",
         M_RED,
         "Fix defects found after release",
         ["Fix auth/login bugs", "Patch SQLite connection errors",
          "Resolve API 500 errors", "Correct twin-status miscalculations"]),
        ("ADAPTIVE",
         M_ORANGE,
         "Adapt to environmental changes",
         ["Node.js & npm upgrades", "Browser compatibility fixes",
          "OS/environment updates", "SQLite driver migration"]),
        ("PERFECTIVE",
         M_GREEN,
         "Improve existing functionality",
         ["Add hydration / meditation metrics", "Enhance dashboard charts",
          "Optimize DB queries (indexes)", "Improve UX responsiveness"]),
        ("PREVENTIVE",
         M_PURPLE,
         "Prevent future failures",
         ["Code refactoring & cleanup", "Add missing test coverage",
          "Update README & API docs", "Dependency vulnerability scans"]),
    ]

    for ax, (title, color, subtitle, items) in zip(axes.flat, types):
        ax.set_facecolor((*color, 0.08))
        ax.set_xlim(0, 10); ax.set_ylim(0, 6.5)
        ax.axis('off')

        # header band
        header = FancyBboxPatch((0, 5.2), 10, 1.3,
                                boxstyle="round,pad=0.05",
                                facecolor=color, edgecolor='none')
        ax.add_patch(header)
        ax.text(5, 5.85, title, ha='center', va='center',
                fontsize=13, fontweight='bold', color='white')
        ax.text(5, 5.35, subtitle, ha='center', va='center',
                fontsize=9, color='white', style='italic')

        for j, item in enumerate(items):
            y = 4.4 - j * 1.1
            dot = plt.Circle((0.6, y + 0.15), 0.18, color=color)
            ax.add_patch(dot)
            ax.text(1.15, y + 0.15, item, va='center', fontsize=10.5, color='#1A1A2E')

    fig.suptitle('Maintenance Types — Boehm Classification',
                 fontsize=15, fontweight='bold', y=1.01, color=tuple(M_DARK))
    plt.savefig('img_types.png', dpi=160, bbox_inches='tight', facecolor='white')
    plt.close()


def make_timeline_diagram():
    """Phased maintenance timeline with swimlane rows."""
    fig, ax = plt.subplots(figsize=(12, 5.5), facecolor='white')
    ax.set_facecolor('white')
    ax.set_xlim(0, 12); ax.set_ylim(-0.5, 5)
    ax.axis('off')

    phases = [
        ("Phase 0\nInitial Release",    0.4, 1.6, M_DARK,   "Core API + Frontend\nDeployed"),
        ("Phase 1\nCorrectiveMaint.", 2.2, 2.2, M_RED,    "Bug Fixes\n& Hotpatches"),
        ("Phase 2\nAdaptive Maint.",  4.8, 2.4, M_ORANGE, "Dependency\nUpdates"),
        ("Phase 3\nPerfective Maint.",7.4, 2.2, M_GREEN,  "Feature\nEnhancements"),
        ("Phase 4\nOngoing Maint.",  10.0, 1.6, M_MID,    "Continuous\nImprovement"),
    ]

    arrow_y = 2.0
    ax.annotate('', xy=(11.8, arrow_y), xytext=(0.1, arrow_y),
                arrowprops=dict(arrowstyle='->', lw=2.5, color='#444'))

    ax.text(0.1, arrow_y - 0.35, 'Project Start', fontsize=9, color='#555', style='italic')
    ax.text(11.4, arrow_y - 0.35, 'Ongoing', fontsize=9, color='#555', style='italic')

    for i, (name, xstart, width, color, desc) in enumerate(phases):
        bar = FancyBboxPatch((xstart, arrow_y - 0.5), width, 1.0,
                             boxstyle="round,pad=0.08",
                             facecolor=color, edgecolor='white', linewidth=2)
        ax.add_patch(bar)
        # Phase label above bar
        lines = name.split('\n')
        ax.text(xstart + width/2, arrow_y + 0.72, lines[0],
                ha='center', fontsize=9, fontweight='bold', color=color)
        if len(lines) > 1:
            ax.text(xstart + width/2, arrow_y + 0.47, lines[1],
                    ha='center', fontsize=8, color='#555')
        # Description inside bar
        ax.text(xstart + width/2, arrow_y, desc,
                ha='center', va='center', fontsize=8.5,
                fontweight='bold', color='white', linespacing=1.4)
        # vertical dashed line marker
        ax.plot([xstart, xstart], [arrow_y - 0.7, arrow_y + 0.2], '--', color='#aaa', lw=0.8)

    ax.text(6, 4.6, 'Maintenance Timeline — Digital Twin Health Tracker',
            ha='center', fontsize=14, fontweight='bold', color=tuple(M_DARK))

    # Legend
    legend_items = [
        (M_RED,    "Corrective"),
        (M_ORANGE, "Adaptive"),
        (M_GREEN,  "Perfective"),
        (M_MID,    "Continuous"),
    ]
    for k, (c, lbl) in enumerate(legend_items):
        px = 1.5 + k * 2.5
        rect = FancyBboxPatch((px, -0.3), 0.35, 0.35, boxstyle="round,pad=0.02",
                               facecolor=c, edgecolor='none')
        ax.add_patch(rect)
        ax.text(px + 0.5, -0.13, lbl, fontsize=9.5, color='#333', va='center')

    plt.tight_layout()
    plt.savefig('img_timeline.png', dpi=160, bbox_inches='tight', facecolor='white')
    plt.close()


def make_architecture_diagram():
    """System component map with maintenance scope highlighted."""
    fig, ax = plt.subplots(figsize=(12, 6.5), facecolor='white')
    ax.set_xlim(0, 12); ax.set_ylim(0, 6)
    ax.axis('off')

    # Layer bands
    layers = [
        (0.1, 4.5, 11.8, 1.2, M_GREY, "Frontend Layer"),
        (0.1, 2.9, 11.8, 1.4, (0.90, 0.95, 1.0), "Application Layer"),
        (0.1, 1.3, 11.8, 1.4, (0.88, 0.96, 0.90), "Data & Test Layer"),
    ]
    for lx, ly, lw, lh, lc, lname in layers:
        ax.add_patch(FancyBboxPatch((lx, ly), lw, lh, boxstyle="round,pad=0.1",
                                   facecolor=lc, edgecolor='#ccc', linewidth=1))
        ax.text(lx + 0.15, ly + lh - 0.22, lname, fontsize=9, color='#666',
                style='italic', fontweight='bold')

    # Components
    components = [
        # Frontend
        (1.5, 4.65, 2.2, 0.8, "index.html\n(Dashboard)", M_LIGHT),
        (4.3, 4.65, 2.2, 0.8, "login.html\n(Auth UI)", M_LIGHT),
        (7.1, 4.65, 2.2, 0.8, "charts.js\napp.js api.js", M_MID),
        # App Layer
        (0.6, 3.05, 2.0, 0.9, "/api/users\n(Users Route)", M_DARK),
        (3.0, 3.05, 2.3, 0.9, "/api/activities\n(Activities Route)", M_DARK),
        (5.7, 3.05, 2.3, 0.9, "/api/twin-status\n(Twin Route)", M_DARK),
        (8.4, 3.05, 2.8, 0.9, "Express Middleware\n(Validate / Error)", M_MID),
        # Data/Test
        (0.6, 1.45, 2.5, 0.9, "health.db\n(SQLite)", M_GREEN),
        (3.5, 1.45, 2.8, 0.9, "API Tests\n(tests/api.test.js)", M_RED),
        (6.7, 1.45, 3.2, 0.9, "UI Tests\n(Playwright / tests-ui/)", M_ORANGE),
    ]
    for cx, cy, cw, ch, label, color in components:
        ax.add_patch(FancyBboxPatch((cx, cy), cw, ch, boxstyle="round,pad=0.08",
                                   facecolor=color, edgecolor='white', linewidth=1.8,
                                   zorder=2))
        ax.text(cx + cw/2, cy + ch/2, label, ha='center', va='center',
                fontsize=8.8, fontweight='bold', color='white', zorder=3, linespacing=1.3)

    # Arrows between layers
    arrow_pairs = [
        (2.6, 4.65, 2.6, 3.95),
        (5.4, 4.65, 5.4, 3.95),
        (1.6, 3.05, 1.6, 2.35),
        (4.15, 3.05, 4.15, 2.35),
        (6.85, 3.05, 6.85, 2.35),
    ]
    for x1, y1, x2, y2 in arrow_pairs:
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', lw=1.6, color='#666'))

    ax.text(6, 5.82, 'System Architecture — Maintenance Scope',
            ha='center', fontsize=14, fontweight='bold', color=tuple(M_DARK))
    plt.tight_layout()
    plt.savefig('img_arch.png', dpi=160, bbox_inches='tight', facecolor='white')
    plt.close()


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 – Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, 'F4F7FB')

    # Full-width dark banner top half
    add_rect(slide, Inches(0), Inches(0), W, Inches(4.2), C_DARK)
    # accent stripe
    add_rect(slide, Inches(0), Inches(4.2), W, Inches(0.12), C_LIGHT)

    # Project title
    add_tb(slide, Inches(0.7), Inches(0.9), Inches(12), Inches(1.1),
           PROJECT, size=38, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Sub-line 1
    add_tb(slide, Inches(0.7), Inches(2.1), Inches(12), Inches(0.5),
           "Maintenance Process Model", size=20, bold=False,
           color=RGBColor(0xB0, 0xC8, 0xF0), align=PP_ALIGN.LEFT)

    # Model badge
    badge = add_rect(slide, Inches(0.7), Inches(2.85), Inches(5.5), Inches(0.7),
                     C_LIGHT, radius=True)
    badge.line.color.rgb = C_LIGHT
    add_tb(slide, Inches(0.8), Inches(2.9), Inches(5.3), Inches(0.6),
           f"Model: {MODEL}", size=18, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Bottom info boxes
    for i, (label, val) in enumerate([("Date", DATE), ("Team", MEMBERS)]):
        bx = Inches(0.7 + i * 4.0)
        add_rect(slide, bx, Inches(4.6), Inches(3.5), Inches(0.7), C_GREY, C_MID, radius=True)
        add_tb(slide, bx + Inches(0.12), Inches(4.65), Inches(1.1), Inches(0.55),
               label, size=11, bold=True, color=C_MID)
        add_tb(slide, bx + Inches(1.2), Inches(4.65), Inches(2.2), Inches(0.55),
               val, size=11, color=C_TEXT)

    # Bottom accent bar
    add_rect(slide, Inches(0), Inches(6.9), W, Inches(0.6), C_MID)
    add_tb(slide, Inches(0.5), Inches(6.93), Inches(12), Inches(0.45),
           "Software Maintenance Engineering  |  Iterative Enhancement Model",
           size=11, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT)


def slide_model_selection(prs):
    """Slide 2 – Model Selection & Justification."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 'F4F7FB')
    slide_header(slide, "Maintenance Process Model Selection",
                 "Project: " + PROJECT)

    # Selected model highlight box
    add_rect(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.75), C_MID, radius=True)
    add_tb(slide, Inches(0.75), Inches(1.68), Inches(12), Inches(0.58),
           f"✔  Selected Model:  {MODEL}", size=20, bold=True,
           color=C_WHITE, align=PP_ALIGN.LEFT)

    # Two columns: Why this model | Models considered
    add_divider(slide, Inches(2.55), left=Inches(0.6), width=Inches(12.13))

    add_tb(slide, Inches(0.6), Inches(2.65), Inches(5.9), Inches(0.35),
           "Why Iterative Enhancement?", size=15, bold=True, color=C_DARK)
    justifications = [
        ("Continuous Evolution",
         "Digital twin health-score algorithm needs periodic recalibration as more user activity data is gathered."),
        ("User-Driven Feedback",
         "Health metrics (steps, sleep, calories, workouts) require iterative UX improvements each cycle."),
        ("Safe Schema Growth",
         "SQLite schema evolves as new activity types are added; iterative cycles enable safe migrations."),
        ("API Enhancement",
         "REST endpoints need backward-compatible extensions without breaking existing clients."),
        ("Ongoing Calibration",
         "Health thresholds and rating logic must be tuned periodically against real-world data."),
    ]
    y = Inches(3.1)
    for title, body in justifications:
        add_tb(slide, Inches(0.7), y, Inches(5.6), Inches(0.28),
               f"▸  {title}", size=12, bold=True, color=C_MID)
        add_tb(slide, Inches(0.8), y + Inches(0.29), Inches(5.5), Inches(0.34),
               body, size=11, color=C_SUB)
        y += Inches(0.71)

    # Right column – other models considered
    add_rect(slide, Inches(7.0), Inches(2.55), Inches(5.73), Inches(4.0), C_GREY, C_LIGHT, radius=True)
    add_tb(slide, Inches(7.12), Inches(2.65), Inches(5.5), Inches(0.35),
           "Other Models Considered", size=15, bold=True, color=C_DARK)

    others = [
        ("Quick Fix Model",
         "Rejected — reactive only, no structured\nphased improvement cycle."),
        ("Boehm Model",
         "Partially adopted — used for maintenance\ntype classification only."),
        ("Osborne Model",
         "Rejected — suited for legacy COBOL-era\nsystems, not a modern Node.js stack."),
        ("Reuse-Oriented Model",
         "Not applicable — system is built from\nscratch, limited reuse scope."),
    ]
    oy = Inches(3.1)
    for name, reason in others:
        add_tb(slide, Inches(7.12), oy, Inches(5.4), Inches(0.3),
               f"✗  {name}", size=12, bold=True, color=RGBColor(0x7F, 0x00, 0x00))
        add_tb(slide, Inches(7.25), oy + Inches(0.3), Inches(5.3), Inches(0.38),
               reason, size=11, color=C_SUB)
        oy += Inches(0.85)


def slide_process_flow(prs):
    """Slide 3 – Iterative Cycle Diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 'F4F7FB')
    slide_header(slide, "Iterative Enhancement Process Flow",
                 "Six-step maintenance cycle applied to the Digital Twin system")
    slide.shapes.add_picture('img_cycle.png',
                             Inches(0.4), Inches(1.6), width=Inches(6.4))

    # Step descriptions on the right
    steps_desc = [
        ("1. Identify Change",   "Bug reports, user feedback, or tech-debt items raise a maintenance need."),
        ("2. Analyze Impact",    "Assess affected modules (routes, models, DB schema, tests)."),
        ("3. Design Solution",   "Draft code changes; update API contracts or DB migration scripts."),
        ("4. Implement & Code",  "Apply changes in feature branch; follow coding standards."),
        ("5. Test & Validate",   "Run api.test.js (Node runner) + Playwright UI suite; check coverage."),
        ("6. Deploy & Monitor",  "Merge to main, restart server, monitor /api/health & logs."),
    ]
    sy = Inches(1.65)
    for title, desc in steps_desc:
        add_tb(slide, Inches(7.0), sy, Inches(5.9), Inches(0.3),
               title, size=13, bold=True, color=C_DARK)
        add_tb(slide, Inches(7.1), sy + Inches(0.30), Inches(5.7), Inches(0.36),
               desc, size=11, color=C_SUB)
        sy += Inches(0.90)


def slide_maintenance_types(prs):
    """Slide 4 – Boehm Classification."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 'F4F7FB')
    slide_header(slide, "Maintenance Types — Boehm Classification",
                 "Every change request is categorised before work begins")
    slide.shapes.add_picture('img_types.png',
                             Inches(0.4), Inches(1.6), width=Inches(12.5))


def slide_activities(prs):
    """Slide 5 – Maintenance Activities table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 'F4F7FB')
    slide_header(slide, "Project-Specific Maintenance Activities",
                 "Mapped to Boehm categories for the Digital Twin Health Tracker")

    # Draw a styled table manually using shapes
    headers = ["Type", "Activity", "Module(s) Affected", "Priority"]
    col_widths = [Inches(1.65), Inches(4.6), Inches(3.8), Inches(2.28)]
    col_x = [Inches(0.5)]
    for cw in col_widths[:-1]:
        col_x.append(col_x[-1] + cw)
    row_h = Inches(0.54)
    start_y = Inches(1.6)

    # Header row
    for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
        add_rect(slide, cx, start_y, cw - Inches(0.03), row_h, C_DARK)
        add_tb(slide, cx + Inches(0.06), start_y + Inches(0.12),
               cw - Inches(0.12), row_h - Inches(0.12),
               hdr, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("CORRECTIVE", C_RED,
         "Fix twin-status rating bug (0-total shows 'good')",
         "src/routes/twin.js",  "High"),
        ("CORRECTIVE", C_RED,
         "Resolve auth missing-password 400 response",
         "src/routes/users.js", "High"),
        ("ADAPTIVE", C_ORANGE,
         "Upgrade better-sqlite3 after Node.js v22 release",
         "package.json, database.js", "Medium"),
        ("ADAPTIVE", C_ORANGE,
         "Ensure Playwright tests pass on latest Chromium",
         "tests-ui/ specs", "Medium"),
        ("PERFECTIVE", C_GREEN,
         "Add weekly trend chart to dashboard",
         "public/js/charts.js", "Medium"),
        ("PERFECTIVE", C_GREEN,
         "Add hydration activity type + seeding data",
         "models/Activity.js, seed.js", "Low"),
        ("PREVENTIVE", C_PURPLE,
         "Increase API test coverage to 80%+",
         "tests/api.test.js", "Medium"),
        ("PREVENTIVE", C_PURPLE,
         "Refactor middleware to reduce code duplication",
         "src/middleware/", "Low"),
    ]

    for ri, (rtype, rcolor, activity, modules, priority) in enumerate(rows):
        ry = start_y + row_h + ri * row_h
        bg = C_GREY if ri % 2 == 0 else C_WHITE
        for cx, cw in zip(col_x, col_widths):
            add_rect(slide, cx, ry, cw - Inches(0.03), row_h - Pt(1), bg)

        # Type cell with color badge
        add_rect(slide, col_x[0] + Inches(0.06), ry + Inches(0.09),
                 col_widths[0] - Inches(0.15), row_h - Inches(0.20),
                 rcolor, radius=True)
        add_tb(slide, col_x[0] + Inches(0.08), ry + Inches(0.13),
               col_widths[0] - Inches(0.19), row_h - Inches(0.26),
               rtype, size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        add_tb(slide, col_x[1] + Inches(0.08), ry + Inches(0.1),
               col_widths[1] - Inches(0.14), row_h - Inches(0.14),
               activity, size=11, color=C_TEXT)
        add_tb(slide, col_x[2] + Inches(0.08), ry + Inches(0.1),
               col_widths[2] - Inches(0.14), row_h - Inches(0.14),
               modules, size=10.5, italic=True, color=C_SUB)

        pri_colors = {"High": C_RED, "Medium": C_ORANGE, "Low": C_GREEN}
        add_tb(slide, col_x[3] + Inches(0.08), ry + Inches(0.1),
               col_widths[3] - Inches(0.14), row_h - Inches(0.14),
               priority, size=11, bold=True, color=pri_colors.get(priority, C_TEXT),
               align=PP_ALIGN.CENTER)


def slide_timeline(prs):
    """Slide 6 – Timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 'F4F7FB')
    slide_header(slide, "Maintenance Timeline & Release Phases",
                 "Planned phases following initial deployment")
    slide.shapes.add_picture('img_timeline.png',
                             Inches(0.3), Inches(1.6), width=Inches(12.7))


def slide_architecture(prs):
    """Slide 7 – Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 'F4F7FB')
    slide_header(slide, "System Architecture — Maintenance Scope",
                 "Components requiring active maintenance attention")
    slide.shapes.add_picture('img_arch.png',
                             Inches(0.3), Inches(1.6), width=Inches(12.7))


def slide_tools(prs):
    """Slide 8 – Tools & Procedures."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 'F4F7FB')
    slide_header(slide, "Maintenance Tools & Procedures",
                 "Standard toolchain used across all maintenance phases")

    tools = [
        (C_DARK,   "Version Control",  "Git",
         "Feature branches, PR reviews, tags per release. Rollback via git revert."),
        (C_RED,    "Bug Tracking",     "GitHub Issues",
         "Issues labelled corrective/adaptive/perfective/preventive before sprint."),
        (C_ORANGE, "API Testing",      "Node Test Runner",
         "tests/api.test.js — run npm test. Target: >80% route coverage."),
        (C_GREEN,  "UI Testing",       "Playwright",
         "tests-ui/ specs (auth, dashboard, activities). Run npm run test:ui."),
        (C_MID,    "Monitoring",       "/api/health + Logs",
         "Health endpoint polled every 60s. Error middleware logs to stdout."),
        (C_PURPLE, "Documentation",    "README + JSDoc",
         "Setup guide, API reference, and inline comments kept up-to-date."),
    ]

    row_h = Inches(1.55)
    for i, (color, category, tool, desc) in enumerate(tools):
        rx = Inches(0.5) + (i % 2) * Inches(6.45)
        ry = Inches(1.6) + (i // 2) * row_h
        card_h = row_h - Inches(0.12)
        # colour left accent bar
        add_rect(slide, rx, ry, Inches(0.1), card_h, color)
        add_rect(slide, rx + Inches(0.12), ry, Inches(6.13), card_h,
                 C_GREY, C_MID, radius=True)
        add_tb(slide, rx + Inches(0.26), ry + Inches(0.12),
               Inches(2.8), Inches(0.38), category, size=14, bold=True, color=C_DARK)
        add_tb(slide, rx + Inches(3.2), ry + Inches(0.12),
               Inches(2.95), Inches(0.38), f"[{tool}]", size=13, bold=True, color=color)
        add_tb(slide, rx + Inches(0.26), ry + Inches(0.58),
               Inches(5.9), Inches(0.72), desc, size=12, color=C_SUB)


def slide_conclusion(prs):
    """Slide 9 – Conclusion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 'F4F7FB')
    slide_header(slide, "Conclusion",
                 "Summary of the Maintenance Process Model for " + PROJECT)

    # Summary paragraph
    para = (
        f"The {PROJECT} adopts the Iterative Enhancement Model as its core maintenance "
        "strategy. This model best fits the project's need for continuous improvement, "
        "regular algorithm recalibration, and user-driven feature evolution while "
        "keeping the system stable and tested throughout its lifecycle."
    )
    add_rect(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.0), C_GREY, C_MID, radius=True)
    add_tb(slide, Inches(0.65), Inches(1.67), Inches(12.0), Inches(0.88),
           para, size=13, color=C_TEXT)

    # Four key takeaways
    takeaways = [
        (C_RED,    "Corrective",  "All reported defects logged as GitHub Issues and triaged before each sprint."),
        (C_ORANGE, "Adaptive",    "Dependencies pinned via package-lock.json; updated every minor sprint."),
        (C_GREEN,  "Perfective",  "Feature roadmap driven by health-domain insights and user feedback."),
        (C_PURPLE, "Preventive",  "80%+ test coverage target + code reviews enforced on every pull request."),
    ]
    for i, (color, title, body) in enumerate(takeaways):
        tx = Inches(0.5) + (i % 2) * Inches(6.42)
        ty = Inches(2.85) + (i // 2) * Inches(1.6)
        add_rect(slide, tx, ty, Inches(6.1), Inches(1.45), color, radius=True)
        add_tb(slide, tx + Inches(0.18), ty + Inches(0.12),
               Inches(5.7), Inches(0.42), title,
               size=16, bold=True, color=C_WHITE)
        add_tb(slide, tx + Inches(0.18), ty + Inches(0.58),
               Inches(5.7), Inches(0.72), body,
               size=12, color=C_WHITE)

    # Footer
    add_rect(slide, Inches(0), Inches(6.9), W, Inches(0.6), C_DARK)
    add_tb(slide, Inches(0.5), Inches(6.93), Inches(12), Inches(0.45),
           f"Maintenance Process Model  |  {PROJECT}  |  {DATE}",
           size=11, color=C_WHITE, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating diagrams …")
    make_iterative_cycle_diagram()
    make_maintenance_types_diagram()
    make_timeline_diagram()
    make_architecture_diagram()

    print("Building presentation …")
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_model_selection(prs)
    slide_process_flow(prs)
    slide_maintenance_types(prs)
    slide_activities(prs)
    slide_timeline(prs)
    slide_architecture(prs)
    slide_tools(prs)
    slide_conclusion(prs)

    out = "Maintenance_Process_Model_Digital_Twin.pptx"
    prs.save(out)
    print(f"Saved: {out}")

    print("Cleaning up temp images …")
    for f in ['img_cycle.png', 'img_types.png', 'img_timeline.png', 'img_arch.png']:
        if os.path.exists(f):
            os.remove(f)
    print("Done!")
