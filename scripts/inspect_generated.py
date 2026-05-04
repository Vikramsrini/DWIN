from pptx import Presentation

prs = Presentation('Maintenance_Process_Model_Digital_Twin.pptx')
for i, slide in enumerate(prs.slides, 1):
    print(f'\n=== SLIDE {i} ===')
    print(f'Layout: {slide.slide_layout.name}')
    for j, shape in enumerate(slide.shapes, 1):
        if hasattr(shape, 'text') and shape.text.strip():
            print(f'  [{shape.shape_type}] {shape.name}: {shape.text[:80]}')
        elif shape.shape_type == 13:
            print(f'  [IMAGE] {shape.name}')
