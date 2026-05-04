from pptx import Presentation
prs = Presentation('Template for Maintenance Process Model.pptx')
for i, slide in enumerate(prs.slides, 1):
    print(f'Slide {i}:')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(f'  - {shape.text[:80]}')
        if shape.has_table:
            table = shape.table
            for r in range(min(3, len(table.rows))):
                row = table.rows[r]
                cells = [cell.text[:30] for cell in row.cells]
                print(f'    Table row {r}: {cells}')
        if shape.has_chart:
            print(f'    [Has chart]')
    print()
